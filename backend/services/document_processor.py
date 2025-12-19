
import os
from pypdf import PdfReader
import docx
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract

# Configure Tesseract path if needed (e.g., specific install location)
# pytesseract.pytesseract.tesseract_cmd = r'/usr/local/bin/tesseract'

class DocumentProcessor:
    @staticmethod
    def extract_text(filepath):
        extracted_text = ""
        filename = os.path.basename(filepath).lower()
        
        try:
            # 1. PDF
            if filename.endswith('.pdf'):
                reader = PdfReader(filepath)
                for page in reader.pages:
                    extracted_text += (page.extract_text() or "") + "\n"
            
            # 2. Word (.docx)
            elif filename.endswith('.docx'):
                doc = docx.Document(filepath)
                for para in doc.paragraphs:
                    extracted_text += para.text + "\n"
            
            # 3. Excel (.xlsx, .xls) & CSV
            elif filename.endswith(('.xlsx', '.xls', '.csv')):
                if filename.endswith('.csv'):
                    df = pd.read_csv(filepath)
                else:
                    df = pd.read_excel(filepath)
                # Convert dataframe to string representation
                extracted_text = df.to_string(index=False)

            # 4. PowerPoint (.pptx)
            elif filename.endswith('.pptx'):
                prs = Presentation(filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"

            # 5. Images (.png, .jpg, .jpeg)
            elif filename.endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp')):
                image = Image.open(filepath)
                extracted_text = pytesseract.image_to_string(image)

            # 6. Text (.txt)
            elif filename.endswith('.txt'):
                 with open(filepath, 'r', encoding='utf-8') as f:
                     extracted_text = f.read()

            else:
                 raise ValueError(f"Unsupported file type: {filename}")
            
            return extracted_text.strip()
        except Exception as e:
            print(f"Error extracting text from {filename}: {e}")
            raise e

    @staticmethod
    def chunk_text(text, chunk_size=1000):
        return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
